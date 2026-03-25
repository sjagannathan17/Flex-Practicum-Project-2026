"""
PowerPoint export — AI-generated content pipeline.

Flow:
  1. Retrieve RAG context via smart_search (multiple topic queries)
  2. Call LLM with structured prompt → JSON with 10 slide specs
  3. Parse JSON response (strip markdown fences, validate)
  4. Render python-pptx slides from parsed data
  5. Return PPTX bytes to caller

The existing Excel pipeline is untouched (separate module).
"""
from __future__ import annotations

import io
import json
import logging
import re
from datetime import datetime
from typing import Optional

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import pptx.oxml.ns as nsmap
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ── Colour palette ─────────────────────────────────────────────────────────
_NAVY         = (26,  39,  68)    # #1a2744 — title slides, table headers
_BLUE         = (59,  130, 246)   # #3b82f6 — accents, Flex highlights
_RED          = (220, 38,  38)    # #dc2626 — HIGH threat / gap boxes
_ORANGE       = (234, 88,  12)    # #ea580c — MEDIUM threat
_GREEN_BADGE  = (22,  163, 74)    # #16a34a — LOW threat / strength boxes
_WHITE        = (255, 255, 255)
_OFF_WHITE    = (248, 250, 252)   # #f8fafc — title bar bg, alt table rows
_LIGHT_BLUE   = (239, 246, 255)   # #eff6ff — footer, callout bg
_LIGHT_ORANGE = (255, 237, 213)   # #ffedd5 — alternating callout bg
_BLUE_TEXT    = (29,  78,  216)   # #1d4ed8 — footer text, Flex callouts
_DARK_GRAY    = (55,  65,  81)    # #374151 — body text
_MID_GRAY     = (100, 116, 139)   # #64748b — captions, secondary text
_DATA_NEEDED  = (156, 163, 175)   # #9ca3af — [Data needed] italic
_FLEX_ROW_BG  = (219, 234, 254)   # #dbeafe — Flex table row highlight

# ── Slide dimensions — widescreen 16:9 ────────────────────────────────────
_SLIDE_W  = Inches(13.33)
_SLIDE_H  = Inches(7.5)
_FOOTER_Y = Inches(7.08)
_FOOTER_H = Inches(0.42)


# ═══════════════════════════════════════════════════════════════════════════
# STEP 1 — RAG CONTEXT RETRIEVAL
# ═══════════════════════════════════════════════════════════════════════════

def _build_rag_context(company: Optional[str], max_chars: int = 12000) -> str:
    """
    Run several topic-focused searches and combine the top passages into a
    single string for the LLM prompt.  Caps total length at max_chars to keep
    the prompt cost predictable.
    """
    try:
        from backend.rag.retriever import smart_search
    except Exception as exc:
        logger.warning("RAG retriever unavailable: %s", exc)
        return "[RAG context unavailable — proceeding without document evidence]"

    # Base queries covering all 10 slide topics
    base_queries = [
        "AI data center revenue competitive strategy Flex Jabil Celestica Benchmark Sanmina",
        "CapEx capital expenditure investment financial performance revenue margin growth",
        "hyperscaler customer relationships Amazon Microsoft Google AI server infrastructure",
        "manufacturing capacity geographic footprint expansion facilities regions",
        "strategic partnerships acquisitions announcements risks tariff supply chain",
        # Flex-specific passes to maximise Flex evidence in context
        "Flex FY2025 FY2026 revenue guidance earnings call management commentary",
        "Flex AI data center segment revenue liquid cooling power modules order pipeline",
        "Flex competitive positioning strategy artificial intelligence infrastructure",
    ]

    # If scoped to one company, add a company-specific pass
    if company and company.lower() != "all":
        base_queries.insert(0, f"{company} competitive position AI strategy earnings revenue")

    seen_ids: set[str] = set()
    passages: list[str] = []
    char_budget = max_chars

    for query in base_queries:
        if char_budget <= 0:
            break
        try:
            result = smart_search(query, n_results=8, use_reranking=False)
            docs = result.get("results", [])
            for doc in docs:
                doc_id = doc.get("id") or doc.get("chunk_id") or ""
                if doc_id and doc_id in seen_ids:
                    continue
                seen_ids.add(doc_id)

                company_label = doc.get("company", "Unknown")
                fy_label      = doc.get("fiscal_year", "")
                filing        = doc.get("filing_type", "")
                content       = (doc.get("parent_content") or doc.get("content", "")).strip()
                if not content:
                    continue

                # Trim each passage to keep context dense
                snippet = content[:600]
                header  = f"[{company_label} | {filing} {fy_label}]\n"
                entry   = header + snippet + "\n\n"

                if len(entry) > char_budget:
                    entry = entry[:char_budget]

                passages.append(entry)
                char_budget -= len(entry)
                if char_budget <= 0:
                    break
        except Exception as exc:
            logger.warning("RAG search failed for query '%s': %s", query[:40], exc)
            continue

    if not passages:
        return "[No document evidence retrieved — use general knowledge of EMS sector]"

    return "".join(passages)


# ═══════════════════════════════════════════════════════════════════════════
# STEP 2 — LLM CALL
# ═══════════════════════════════════════════════════════════════════════════

_SYSTEM_PROMPT = """\
You are a senior competitive intelligence analyst at Flex preparing an executive \
PowerPoint briefing for Flex's strategy leadership team.

Context: Flex competes with Jabil, Celestica, Benchmark, and Sanmina in EMS. \
The primary strategic question is: How aggressively are competitors positioning \
in AI/data center infrastructure, and what does this mean for Flex's competitive \
position and order pipeline?

CRITICAL DATA RULES — follow exactly:
1. Every slide MUST include data or directional commentary for ALL five companies: \
Flex, Jabil, Celestica, Benchmark, Sanmina. Never omit any company.
2. For Flex specifically: if an exact figure is not in the retrieved context, \
you MUST write the directional fallback in this format — \
'Flex: [specific metric] not disclosed in public filings — management commentary suggests [directional signal]'. \
Do NOT leave Flex rows blank. Do NOT write [Data needed] for Flex.
3. Do NOT write [Data needed] for any company. Replace with the best available \
directional signal drawn from the retrieved evidence or sector knowledge.
4. Use the anchor data provided in the user prompt as ground truth for competitive rankings.
5. For Slide 3, rank companies by AI/DC revenue as % of total revenue, \
not by product launches or announcements.\
"""

_USER_TEMPLATE = """\
Retrieved evidence:
{rag_context}

── ANCHOR DATA (treat as ground truth) ──────────────────────────────────────
AI/DC Revenue Mix (% of total revenue, most recent fiscal year):
  Celestica:  ~45%  (HPS segment, accelerating — fastest growing peer)
  Jabil:      ~35%  (cloud & infrastructure, mix accelerating)
  Flex:       ~20–25% (AI Infrastructure & Cloud segment, growing)
  Sanmina:    ~15–20% (cloud & communications mix, growing)
  Benchmark:  ~10%  (higher-value compute, smallest AI/DC exposure)

Flex specific signals (use these to fill Flex rows; supplement with retrieved evidence):
  - FY2025 revenue guidance: management has guided for growth in AI Infrastructure & Cloud
  - Flex liquid cooling: announced active liquid cooling and power module capabilities for AI servers
  - Flex order pipeline: management commentary describes "increased order activity" in AI/DC segment
  - Flex FY2026: not yet guided publicly — directional signal is continued AI/DC mix expansion
  - Flex positioning: "design-to-manufacturing" differentiation vs pure EMS peers; Sketch-to-Scale model
  - When exact Flex figures are absent from evidence, use the format:
    'Flex: [metric] not disclosed in public filings — management commentary suggests [directional signal]'

Competitive ranking rule for Slide 3: rank by AI/DC revenue as % of total — use the anchor percentages above.
─────────────────────────────────────────────────────────────────────────────

Generate a structured 10-slide executive presentation{scope_note}.

Requirements:
- Every slide has ONE clear insight for Flex leadership
- Every slide MUST name all five companies: Flex, Jabil, Celestica, Benchmark, Sanmina
- Compare companies directly, never summarize in isolation
- Use specific numbers and company names in every bullet
- Max 12 words per bullet
- Frame everything from Flex's perspective
- For any Flex data not in retrieved evidence, use the directional fallback format above — never leave Flex blank
- Never use: "significant", "notable", "key opportunity", "strong performance", "important", "various", "several"
- Never write "[Data needed]" — always substitute a directional signal

Return ONLY valid JSON, no markdown, no explanation:

{{
  "report_title": "Flex Competitive Intelligence Brief",
  "generated_date": "{today}",
  "executive_summary_one_liner": "[single most important insight]",
  "slides": [
    {{
      "slide_number": 1,
      "slide_title": "Executive Summary",
      "slide_insight": "[core strategic implication in one sentence]",
      "bullets": ["bullet 1", "bullet 2", "bullet 3", "bullet 4", "bullet 5"],
      "flex_action": "[one concrete action starting with a verb]"
    }},
    {{
      "slide_number": 2,
      "slide_title": "AI Infrastructure Demand — Market Context",
      "slide_insight": "[why timing matters for Flex now]",
      "bullets": ["bullet 1", "bullet 2", "bullet 3", "bullet 4", "bullet 5"],
      "flex_action": "[one concrete action]"
    }},
    {{
      "slide_number": 3,
      "slide_title": "AI/DC Competitive Ranking",
      "slide_insight": "Celestica leads at ~45% AI/DC mix; Flex at ~20-25% must close gap or cede hyperscaler share",
      "bullets": [
        "Rank 1: Celestica — ~45% AI/DC revenue mix, HPS segment fastest-growing",
        "Rank 2: Jabil — ~35% AI/DC mix, cloud & infrastructure accelerating",
        "Rank 3: Flex — ~20-25% AI/DC mix, AI Infrastructure & Cloud growing",
        "Rank 4: Sanmina — ~15-20% mix, cloud & communications expanding",
        "Rank 5: Benchmark — ~10% mix, highest-value compute, smallest AI/DC exposure",
        "Fastest accelerating: Celestica — HPS segment doubling year-over-year"
      ],
      "flex_action": "[one concrete action to close the gap to Celestica/Jabil]"
    }},
    {{
      "slide_number": 4,
      "slide_title": "Capacity & Geographic Footprint",
      "slide_insight": "[who has capacity where demand is growing]",
      "bullets": [
        "Flex: [specific footprint signal from evidence or directional fallback]",
        "Celestica: [specific]",
        "Jabil: [specific]",
        "Sanmina: [specific]",
        "Benchmark: [specific]"
      ],
      "flex_action": "[one concrete action]"
    }},
    {{
      "slide_number": 5,
      "slide_title": "Financial Performance & Growth Momentum",
      "slide_insight": "[revenue and margin comparison insight]",
      "bullets": [
        "Flex: [revenue/CapEx figure or directional fallback]",
        "Celestica: [specific]",
        "Jabil: [specific]",
        "Sanmina: [specific]",
        "Benchmark: [specific]"
      ],
      "flex_action": "[one concrete action]"
    }},
    {{
      "slide_number": 6,
      "slide_title": "Hyperscaler Customer Intelligence",
      "slide_insight": "[who has strongest hyperscaler relationships]",
      "bullets": [
        "Flex: [hyperscaler signal from evidence or directional fallback]",
        "Celestica: [specific]",
        "Jabil: [specific]",
        "Sanmina: [specific]",
        "Benchmark: [specific]"
      ],
      "flex_action": "[one concrete action]"
    }},
    {{
      "slide_number": 7,
      "slide_title": "Strategic Moves & Announcements",
      "slide_insight": "[most impactful recent move for Flex]",
      "bullets": [
        "Flex: [strategic move or directional fallback — liquid cooling, power modules, Sketch-to-Scale]",
        "Celestica: [specific]",
        "Jabil: [specific]",
        "Sanmina: [specific]",
        "Benchmark: [specific]"
      ],
      "flex_action": "[one concrete action]"
    }},
    {{
      "slide_number": 8,
      "slide_title": "Risk Landscape",
      "slide_insight": "[risk most directly threatening Flex]",
      "bullets": ["bullet 1", "bullet 2", "bullet 3", "bullet 4", "bullet 5"],
      "flex_action": "[one concrete action]"
    }},
    {{
      "slide_number": 9,
      "slide_title": "Flex Competitive Position",
      "slide_insight": "[honest strengths and gaps assessment]",
      "bullets": [
        "Strength 1: Flex design-to-manufacturing Sketch-to-Scale differentiation",
        "Strength 2: [specific from evidence]",
        "Gap 1: AI/DC revenue mix ~20-25% vs Celestica ~45% — 20pp gap to close",
        "Gap 2: [specific vs Jabil or Celestica]",
        "Net trajectory: improving — AI Infrastructure & Cloud segment growing"
      ],
      "flex_action": "[highest priority gap to close]"
    }},
    {{
      "slide_number": 10,
      "slide_title": "Recommended Next Steps",
      "slide_insight": "[three actions in next 90 days]",
      "bullets": [
        "Action 1 [30 days]: [specific and measurable]",
        "Action 2 [60 days]: [specific and measurable]",
        "Action 3 [90 days]: [specific and measurable]",
        "Success metric: [what good looks like in 6 months]",
        "Key assumption to validate: [before committing resources]"
      ],
      "flex_action": "[owner and timeline for follow-up]"
    }}
  ]
}}\
"""


def _call_llm_for_report(rag_context: str, company: Optional[str]) -> str:
    """Return the raw LLM response string."""
    from backend.core.llm_client import llm_complete

    scope_note = ""
    if company and company.lower() not in ("all", "comparison"):
        scope_note = f" focused on {company.title()} vs peers"

    user_msg = _USER_TEMPLATE.format(
        rag_context=rag_context,
        scope_note=scope_note,
        today=datetime.now().strftime("%B %d, %Y"),
    )

    return llm_complete(
        messages=[{"role": "user", "content": user_msg}],
        system=_SYSTEM_PROMPT,
        model_key="main",
        max_tokens=6000,
        stream=False,
    )


# ═══════════════════════════════════════════════════════════════════════════
# STEP 3 — PARSE JSON RESPONSE
# ═══════════════════════════════════════════════════════════════════════════

def _parse_llm_response(raw: str) -> dict:
    """
    Strip markdown fences (if any) and parse JSON.
    Raises ValueError with descriptive message on failure.
    """
    text = raw.strip()

    # Strip ```json ... ``` or ``` ... ``` fences
    fenced = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", text)
    if fenced:
        text = fenced.group(1).strip()
    else:
        # Grab first {...} block as a fallback
        m = re.search(r"\{[\s\S]*\}", text)
        if m:
            text = m.group(0)

    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        logger.error("LLM JSON parse failed. Raw response head:\n%s", raw[:500])
        raise ValueError(
            "Report generation failed — AI response could not be parsed. "
            "Please try again."
        ) from exc

    # Basic schema validation
    if "slides" not in data or not isinstance(data["slides"], list):
        raise ValueError(
            "Report generation failed — AI response missing slide data. "
            "Please try again."
        )

    return data


# ═══════════════════════════════════════════════════════════════════════════
# STEP 4 — PYTHON-PPTX RENDERING  (executive visual redesign)
# ═══════════════════════════════════════════════════════════════════════════

def _rgb(r: int, g: int, b: int) -> "RGBColor":
    return RGBColor(r, g, b)


# ── Layout constants ────────────────────────────────────────────────────────
_TITLE_BAR_H = Inches(1.45)
_CONTENT_TOP = Inches(1.52)
_COL_L       = Inches(0.40)          # default left margin
_COL2_L      = Inches(8.95)          # right-column start
_COL2_W      = Inches(4.00)          # right-column width
_CONTENT_H   = _FOOTER_Y - _CONTENT_TOP - Inches(0.08)

# Hardcoded supplementary callout data for specific LLM slide numbers.
# Keys are the LLM slide_number values (1-10); values are lists of
# (metric, label, bg_colour) triples used instead of _extract_metrics().
_SLIDE_CALLOUT_DATA: dict[int, list[tuple]] = {
    2: [  # AI Infrastructure Demand
        ("$675B",  "2026 Big Five CapEx",         _LIGHT_BLUE),
        ("+128%",  "Average YoY Growth",           _LIGHT_ORANGE),
        ("7GW",    "Planned AI Capacity",           _LIGHT_BLUE),
    ],
    4: [  # Capacity & Footprint
        ("30",     "Countries — Flex Footprint",   _LIGHT_BLUE),
        ("160K",   "Flex Employees",               _LIGHT_ORANGE),
        ("7",      "Flex-Competitor Overlap Zones", _LIGHT_BLUE),
    ],
    6: [  # Hyperscaler Intelligence
        ("$200B",  "AWS 2026 CapEx",               _LIGHT_BLUE),
        ("$180B",  "Alphabet 2026 CapEx",          _LIGHT_ORANGE),
        ("+212%",  "Meta YoY CapEx Growth",        _LIGHT_BLUE),
    ],
}

# Slide numbers (LLM) whose right column uses insight-phrase boxes
# instead of numeric stat callouts.
_INSIGHT_PHRASE_SLIDES = {7, 8}


# ── Low-level primitives ────────────────────────────────────────────────────

def _rect(slide, left, top, width, height, fill, line=None):
    """Filled rectangle; no line by default."""
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(*fill)
    if line:
        s.line.color.rgb = _rgb(*line)
    else:
        s.line.fill.background()
    return s


def _tbox(slide, text, left, top, width, height, *,
          size, bold=False, italic=False,
          color=_DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True):
    """Textbox with a single formatted run."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = _rgb(*color)
    return box


def _set_cell(cell, text, *, size, bold=False, color=_DARK_GRAY,
              align=PP_ALIGN.CENTER):
    """Set text + formatting on a fresh table cell."""
    cell.margin_left   = Inches(0.08)
    cell.margin_right  = Inches(0.08)
    cell.margin_top    = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = _rgb(*color)
    try:
        from pptx.enum.text import MSO_ANCHOR
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass


def _stat_box(slide, metric, label, left, top, width, height,
              bg=_LIGHT_BLUE):
    """Stat callout: large metric + small label in a coloured box."""
    _rect(slide, left, top, width, height, fill=bg, line=_BLUE)
    # Large metric
    _tbox(slide, metric,
          left + Inches(0.12), top + Inches(0.12),
          width - Inches(0.24), Inches(0.72),
          size=34, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)
    # Small label
    _tbox(slide, label,
          left + Inches(0.12), top + Inches(0.88),
          width - Inches(0.24), Inches(0.45),
          size=11, color=_MID_GRAY, align=PP_ALIGN.CENTER)


def _footer_strip(slide, flex_action):
    """Light-blue footer: Flex Action left, confidential right."""
    _rect(slide, Inches(0), _FOOTER_Y, _SLIDE_W, _FOOTER_H, fill=_LIGHT_BLUE)

    box = slide.shapes.add_textbox(
        Inches(0.25), _FOOTER_Y + Inches(0.06), Inches(9.8), Inches(0.3))
    tf = box.text_frame
    p  = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text            = "⚡ Flex Action: "
    r1.font.bold       = True
    r1.font.size       = Pt(10)
    r1.font.color.rgb  = _rgb(*_BLUE_TEXT)
    r2 = p.add_run()
    r2.text            = flex_action
    r2.font.size       = Pt(10)
    r2.font.color.rgb  = _rgb(*_BLUE_TEXT)

    _tbox(slide,
          f"EMS Pulse Intelligence · Confidential · {datetime.now().strftime('%B %Y')}",
          Inches(10.1), _FOOTER_Y + Inches(0.06), Inches(3.0), Inches(0.3),
          size=9, color=_MID_GRAY, align=PP_ALIGN.RIGHT)


def _title_bar(slide, num, title, insight):
    """Off-white title bar with blue left strip, bold title, italic insight."""
    _rect(slide, Inches(0), Inches(0), _SLIDE_W, _TITLE_BAR_H,
          fill=_OFF_WHITE, line=(226, 232, 240))
    _rect(slide, Inches(0), Inches(0), Inches(0.12), _TITLE_BAR_H, fill=_BLUE)

    _tbox(slide, f"{num:02d}  {title}",
          Inches(0.25), Inches(0.10), Inches(13.0), Inches(0.68),
          size=24, bold=True, color=_NAVY)

    if insight:
        _tbox(slide, insight,
              Inches(0.25), Inches(0.78), Inches(13.0), Inches(0.58),
              size=13, italic=True, color=_DARK_GRAY)


def _bullet_col(slide, bullets, left, top, width, height):
    """Render ▸ bullet list; colours [Data needed] entries differently."""
    if not bullets:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        is_dn = "[data needed" in text.lower()
        run = p.add_run()
        run.text           = f"▸  {text}"
        run.font.size      = Pt(15)
        run.font.italic    = is_dn
        run.font.color.rgb = _rgb(*_DATA_NEEDED) if is_dn else _rgb(*_DARK_GRAY)


def _extract_metrics(bullets):
    """Pull up to 3 (metric, label) pairs from bullets for stat callouts."""
    found = []
    pat = re.compile(
        r'([+\-~]?\$?[\d,.]+(?:\s?[BMK%]|\s?billion|\s?million)?)',
        re.IGNORECASE)
    for bullet in bullets:
        m = pat.search(bullet)
        if m:
            metric = m.group(1).strip()
            label  = re.sub(r'\s+', ' ',
                            bullet.replace(metric, '').strip().lstrip('▸:— '))[:42]
            if metric and label:
                found.append((metric, label))
        if len(found) >= 3:
            break
    return found


# ── TITLE SLIDE ─────────────────────────────────────────────────────────────

def _add_title_slide(prs, data):
    """Dark-navy cover: EMS Pulse badge, large title, accent line, date."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _rect(slide, Inches(0), Inches(0), _SLIDE_W, _SLIDE_H, fill=_NAVY)

    # EMS Pulse badge — top-left
    _tbox(slide, "EMS Pulse",
          Inches(0.45), Inches(0.38), Inches(2.8), Inches(0.55),
          size=18, bold=True, color=_BLUE)

    # Report title — centred
    _tbox(slide, data.get("report_title", "Flex Competitive Intelligence Brief"),
          Inches(1.0), Inches(2.25), Inches(11.33), Inches(1.6),
          size=40, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    # Thin accent line below title
    _rect(slide, Inches(1.5), Inches(4.0), Inches(10.33), Inches(0.05), fill=_BLUE)

    # Date
    _tbox(slide, data.get("generated_date", datetime.now().strftime("%B %d, %Y")),
          Inches(1.0), Inches(4.18), Inches(11.33), Inches(0.55),
          size=18, color=(148, 163, 184), align=PP_ALIGN.CENTER)

    # Confidential footer
    _tbox(slide, "Confidential — Flex Internal Use Only",
          Inches(1.0), Inches(6.9), Inches(11.33), Inches(0.42),
          size=11, color=_MID_GRAY, align=PP_ALIGN.CENTER)


# ── EXECUTIVE SUMMARY (Slide 1) ─────────────────────────────────────────────

def _add_exec_summary_slide(prs, slide_data, exec_liner):
    """
    2-column layout:
      Left 60% — highlight banner + large italic quote + bullets
      Right 40% — 3 stat callout boxes with real anchor metrics
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    num     = slide_data.get("slide_number", 1)
    title   = slide_data.get("slide_title", "Executive Summary")
    insight = slide_data.get("slide_insight", "")
    bullets = slide_data.get("bullets", [])
    action  = slide_data.get("flex_action", "")

    _title_bar(slide, num, title, "")   # insight shown in body

    left_w = Inches(7.65)

    # Highlight banner — max 0.70" for single-line, 0.25" padding each side
    if exec_liner:
        liner_lines = max(1, (len(exec_liner) + 79) // 80)
        banner_h = Inches(0.70) if liner_lines == 1 \
                   else Inches(min(0.70 + (liner_lines - 1) * 0.22, 1.0))
        _rect(slide, _COL_L, _CONTENT_TOP, left_w, banner_h,
              fill=_LIGHT_BLUE, line=_BLUE)
        _tbox(slide, exec_liner,
              _COL_L + Inches(0.15), _CONTENT_TOP + Inches(0.09),
              left_w - Inches(0.3), banner_h - Inches(0.14),
              size=13, bold=True, color=_BLUE_TEXT, align=PP_ALIGN.CENTER)

    # Large italic insight quote
    q_top = _CONTENT_TOP + (banner_h + Inches(0.10) if exec_liner else 0)
    if insight:
        _tbox(slide, f'"{insight}"',
              _COL_L, q_top, left_w, Inches(0.95),
              size=18, italic=True, color=_DARK_GRAY)

    # Bullets
    _bullet_col(slide, bullets[:5],
                _COL_L, q_top + Inches(1.02), left_w, Inches(3.8))

    # Right column — 3 stat callouts (fixed anchor metrics)
    right_l = Inches(8.30)
    right_w = Inches(4.65)
    stat_h  = Inches(1.50)
    gap     = Inches(0.18)

    for i, (metric, label, bg) in enumerate([
        ("+35%",  "Jabil FY25 EPS Growth",      _LIGHT_BLUE),
        ("+68%",  "Celestica Revenue Growth",   _LIGHT_ORANGE),
        ("~45%",  "Celestica AI/DC Revenue Mix", _LIGHT_BLUE),
    ]):
        _stat_box(slide, metric, label,
                  right_l, _CONTENT_TOP + i * (stat_h + gap),
                  right_w, stat_h, bg=bg)

    _footer_strip(slide, action)


# ── COMPETITIVE RANKING (Slide 3) ───────────────────────────────────────────

_RANK_ROWS = [
    ("🥇 1", "Celestica", "~45%", "+68%", "🔴 HIGH",    _RED),
    ("🥈 2", "Jabil",     "~35%", "+52%", "🔴 HIGH",    _RED),
    ("  3",  "Flex",      "~22%", "+45%", "—  Baseline", _BLUE),
    ("  4",  "Sanmina",   "~18%", "+30%", "🟠 MID",     _ORANGE),
    ("  5",  "Benchmark", "~10%", "+25%", "🟢 LOW",     _GREEN_BADGE),
]
_RANK_HDRS = ["Rank", "Company", "AI/DC Mix", "YoY Growth", "Threat to Flex"]


def _add_ranking_slide(prs, slide_data):
    """
    Table layout ranked by AI/DC revenue %, not announcements.
    Flex row is highlighted with light-blue bg + left blue accent bar.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    num    = slide_data.get("slide_number", 3)
    title  = slide_data.get("slide_title", "AI/DC Competitive Ranking")
    insight = slide_data.get("slide_insight", "")
    action  = slide_data.get("flex_action", "")

    _title_bar(slide, num, title, insight)

    tbl_l = Inches(0.35)
    tbl_t = Inches(1.58)
    tbl_w = Inches(12.63)
    tbl_h = Inches(5.30)

    shape = slide.shapes.add_table(
        1 + len(_RANK_ROWS), len(_RANK_HDRS), tbl_l, tbl_t, tbl_w, tbl_h)
    tbl = shape.table

    for i, w in enumerate([Inches(1.1), Inches(2.55), Inches(2.3),
                            Inches(2.3),  Inches(4.38)]):
        tbl.columns[i].width = w

    # Header row
    tbl.rows[0].height = Inches(0.60)
    for c, hdr in enumerate(_RANK_HDRS):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(*_NAVY)
        _set_cell(cell, hdr, size=13, bold=True, color=_WHITE)

    # Data rows
    row_h = Inches(0.84)
    for r, (rank, company, mix, yoy, threat, t_color) in enumerate(_RANK_ROWS):
        ri = r + 1
        tbl.rows[ri].height = row_h
        is_flex = (company == "Flex")
        bg = _FLEX_ROW_BG if is_flex else (_OFF_WHITE if r % 2 == 0 else _WHITE)

        for c, val in enumerate([rank, company, mix, yoy, threat]):
            cell = tbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(*bg)
            txt_color = (t_color if c == 4 and not is_flex
                         else _NAVY if is_flex else _DARK_GRAY)
            _set_cell(cell, val, size=14, bold=is_flex, color=txt_color)

    # Blue left-edge accent bar overlaid on the Flex row (row index 3)
    flex_row_top = tbl_t + Inches(0.60) + 2 * row_h
    _rect(slide, tbl_l, flex_row_top, Inches(0.09), row_h, fill=_BLUE)

    _footer_strip(slide, action)


# ── FINANCIAL PERFORMANCE (Slide 5) ─────────────────────────────────────────

_FIN_ROWS = [
    ("Jabil",     "$28.9B (FY24)", "+35% EPS",  "↑ Expanding"),
    ("Celestica", "~$8.5B",        "+68%",       "↑ Expanding"),
    ("Flex",      "~$26B (FY25e)", "~+10–15%",   "→ Stable"),
    ("Sanmina",   "~$7.9B",        "+30% NI",    "→ Stable"),
    ("Benchmark", "~$2.2B",        "+25% NI",    "↓ Cautious"),
]
_FIN_HDRS = ["Company", "Revenue", "YoY Growth", "Margin Trend"]


def _add_financial_slide(prs, slide_data):
    """
    Left: revenue comparison table.
    Right: 3 large stat callout boxes.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    num    = slide_data.get("slide_number", 5)
    title  = slide_data.get("slide_title", "Financial Performance & Growth Momentum")
    insight = slide_data.get("slide_insight", "")
    action  = slide_data.get("flex_action", "")

    _title_bar(slide, num, title, insight)

    # ── Left table ────────────────────────────────────────────────────────
    tbl_l = Inches(0.35)
    tbl_w = Inches(8.1)

    shape = slide.shapes.add_table(
        1 + len(_FIN_ROWS), len(_FIN_HDRS),
        tbl_l, _CONTENT_TOP, tbl_w, Inches(5.40))
    tbl = shape.table

    for i, w in enumerate([Inches(2.1), Inches(2.2), Inches(2.0), Inches(1.8)]):
        tbl.columns[i].width = w

    tbl.rows[0].height = Inches(0.55)
    for c, hdr in enumerate(_FIN_HDRS):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(*_NAVY)
        _set_cell(cell, hdr, size=13, bold=True, color=_WHITE)

    row_h = Inches(0.97)
    for r, (company, rev, yoy, margin) in enumerate(_FIN_ROWS):
        ri = r + 1
        tbl.rows[ri].height = row_h
        is_flex = (company == "Flex")
        bg = _FLEX_ROW_BG if is_flex else (_OFF_WHITE if r % 2 == 0 else _WHITE)

        for c, val in enumerate([company, rev, yoy, margin]):
            cell = tbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(*bg)
            _set_cell(cell, val, size=13, bold=is_flex,
                      color=_NAVY if is_flex else _DARK_GRAY,
                      align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

    # ── Right stat callouts ───────────────────────────────────────────────
    right_l = Inches(8.72)
    right_w = Inches(4.28)
    stat_h  = Inches(1.55)
    gap     = Inches(0.18)

    for i, (metric, label, bg) in enumerate([
        ("$7.8B",  "Jabil Q3 FY25 Revenue",  _LIGHT_BLUE),
        ("+35%",   "Jabil FY25 EPS Growth",   _LIGHT_ORANGE),
        ("2.6%",   "Benchmark CapEx/Revenue", _LIGHT_BLUE),
    ]):
        _stat_box(slide, metric, label,
                  right_l, _CONTENT_TOP + i * (stat_h + gap),
                  right_w, stat_h, bg=bg)

    _footer_strip(slide, action)


# ── FLEX COMPETITIVE POSITION (Slide 9) ─────────────────────────────────────

def _add_flex_position_slide(prs, slide_data):
    """
    2×2 grid: green strength boxes (top) + red gap boxes (bottom).
    Below grid: progress bar showing competitive trajectory.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    num     = slide_data.get("slide_number", 9)
    title   = slide_data.get("slide_title", "Flex Competitive Position")
    insight = slide_data.get("slide_insight", "")
    bullets = slide_data.get("bullets", [])
    action  = slide_data.get("flex_action", "")

    _title_bar(slide, num, title, insight)

    # Parse bullet text into buckets
    strengths   = [b for b in bullets if "strength" in b.lower()]
    gaps        = [b for b in bullets if "gap" in b.lower()]
    traj_bullet = next((b for b in bullets if "trajectory" in b.lower()), "")

    def _strip(t):
        return re.sub(
            r'^(Strength\s*\d+|Gap\s*\d+|Net trajectory)[:\s]*', '',
            t, flags=re.IGNORECASE).strip()

    s1 = _strip(strengths[0]) if len(strengths) > 0 \
        else "Engineering expertise: liquid cooling & power modules"
    s2 = _strip(strengths[1]) if len(strengths) > 1 \
        else "Diversified: 160K employees, 30 countries, Sketch-to-Scale"
    g1 = _strip(gaps[0])      if len(gaps)      > 0 \
        else "AI/DC revenue mix ~22% vs Celestica ~45% — 23pp gap"
    g2 = _strip(gaps[1])      if len(gaps)      > 1 \
        else "Hyperscaler direct qualification vs Jabil and Celestica"

    # Grid geometry
    grid_l  = Inches(0.35)
    grid_t  = _CONTENT_TOP
    cell_w  = Inches(6.18)
    cell_h  = Inches(2.22)
    gap_px  = Inches(0.22)

    for header, body, accent, col, row in [
        ("✅ STRENGTH 1", s1, _GREEN_BADGE, 0, 0),
        ("✅ STRENGTH 2", s2, _GREEN_BADGE, 1, 0),
        ("⚠️  GAP 1",     g1, _RED,         0, 1),
        ("⚠️  GAP 2",     g2, _RED,         1, 1),
    ]:
        cx = grid_l + col * (cell_w + gap_px)
        cy = grid_t  + row * (cell_h + gap_px)

        # Light tinted background
        bg_tint = (240, 253, 244) if accent == _GREEN_BADGE else (254, 242, 242)
        _rect(slide, cx, cy, cell_w, cell_h, fill=bg_tint, line=accent)

        # Coloured header strip
        _rect(slide, cx, cy, cell_w, Inches(0.46), fill=accent)

        # Header label in strip
        _tbox(slide, header,
              cx + Inches(0.15), cy + Inches(0.06),
              cell_w - Inches(0.3), Inches(0.36),
              size=13, bold=True, color=_WHITE)

        # Body text
        _tbox(slide, body,
              cx + Inches(0.18), cy + Inches(0.56),
              cell_w - Inches(0.36), cell_h - Inches(0.68),
              size=14, color=_DARK_GRAY)

    # ── Trajectory progress bar ───────────────────────────────────────────
    # Fixed 10" bar centred on slide; 3 colour-coded zones.
    bar_t   = grid_t + 2 * (cell_h + gap_px) + Inches(0.15)
    bar_l   = Inches(1.67)   # (13.33 - 10) / 2 — centred
    bar_w   = Inches(10)
    track_h = Inches(0.22)

    traj_lc = traj_bullet.lower()
    if "accelerat" in traj_lc or "improving" in traj_lc:
        dot_offset = int(bar_w * 0.72)   # 72% — well inside green zone
    elif "declin" in traj_lc:
        dot_offset = int(bar_w * 0.15)
    else:
        dot_offset = int(bar_w * 0.40)

    # Three coloured track segments: red | amber | green
    _rect(slide, bar_l,              bar_t + Inches(0.10),
          Inches(2), track_h, fill=(220, 38, 38))
    _rect(slide, bar_l + Inches(2),  bar_t + Inches(0.10),
          Inches(4), track_h, fill=(234, 179, 8))
    _rect(slide, bar_l + Inches(6),  bar_t + Inches(0.10),
          Inches(4), track_h, fill=(22, 163, 74))

    # Filled blue circle marker
    dot_d = Inches(0.20)
    circ = slide.shapes.add_shape(
        9,  # msoShapeOval
        bar_l + dot_offset - dot_d // 2,
        bar_t + Inches(0.01),
        dot_d, dot_d)
    circ.fill.solid()
    circ.fill.fore_color.rgb = _rgb(*_BLUE)
    circ.line.color.rgb      = _rgb(*_WHITE)

    # "← Flex" label directly below marker (10pt)
    lbl_w = Inches(0.90)
    _tbox(slide, "← Flex",
          bar_l + dot_offset - lbl_w // 2,
          bar_t + Inches(0.30),
          lbl_w, Inches(0.26),
          size=10, bold=True, color=_BLUE_TEXT, align=PP_ALIGN.CENTER)

    # Section labels at 15% / 45% / 80% of bar width (10pt, centred on position)
    lbl_w2 = Inches(1.60)
    for pct, text in [(0.15, "Declining"), (0.45, "Stable"), (0.80, "Accelerating")]:
        cx = bar_l + int(bar_w * pct)
        _tbox(slide, text,
              cx - lbl_w2 // 2, bar_t + Inches(0.56),
              lbl_w2, Inches(0.26),
              size=10, color=_MID_GRAY, align=PP_ALIGN.CENTER)

    _footer_strip(slide, action)


# ── DEFAULT SLIDE (all other slide numbers) ──────────────────────────────────

def _split_insight_phrases(insight: str, max_phrases: int = 3) -> list[str]:
    """Split an insight sentence into 2-3 short display phrases."""
    text = (insight or "").strip()
    if not text:
        return ["Strategic insight unavailable"]
    # Try semicolons first (clearest phrase boundary)
    parts = [p.strip() for p in text.split(";") if p.strip()]
    if len(parts) >= 2:
        return parts[:max_phrases]
    # Try em-dash / en-dash separators
    parts = [p.strip() for p in re.split(r'\s+[—–\-]{1,2}\s+', text) if p.strip()]
    if len(parts) >= 2:
        return parts[:max_phrases]
    # Try splitting on ", " yielding at least 2 parts
    parts = [p.strip() for p in text.split(", ") if p.strip()]
    if len(parts) >= 2:
        # Group into at most max_phrases chunks
        chunk = max(1, len(parts) // max_phrases)
        grouped = [", ".join(parts[i:i+chunk]) for i in range(0, len(parts), chunk)]
        return grouped[:max_phrases]
    # Return as single phrase (truncated to 120 chars)
    return [text[:120]]


def _add_default_slide(prs, slide_data):
    """
    2-column layout:
      Left 65%  — ▸ bullet list (up to 6 items)
      Right 35% — supplementary callout boxes (hardcoded per slide number,
                  insight-phrases for strategic/risk slides, else key-insight box)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    num     = slide_data.get("slide_number", 0)
    title   = slide_data.get("slide_title", "")
    insight = slide_data.get("slide_insight", "")
    bullets = slide_data.get("bullets", [])
    action  = slide_data.get("flex_action", "")

    _title_bar(slide, num, title, insight)

    left_w = Inches(8.32)
    _bullet_col(slide, bullets[:6], _COL_L, _CONTENT_TOP, left_w, _CONTENT_H)

    stat_h = Inches(1.55)
    gap    = Inches(0.20)

    if num in _SLIDE_CALLOUT_DATA:
        # Hardcoded supplementary stats — never derived from bullet text
        for i, (metric, label, bg) in enumerate(_SLIDE_CALLOUT_DATA[num]):
            _stat_box(slide, metric, label,
                      _COL2_L, _CONTENT_TOP + i * (stat_h + gap),
                      _COL2_W, stat_h, bg=bg)

    elif num == 7:
        # Slide 8 — Strategic Moves: "Key Competitive Moves" list
        panel_h = Inches(4.80)
        _rect(slide, _COL2_L, _CONTENT_TOP, _COL2_W, panel_h,
              fill=_LIGHT_BLUE, line=_BLUE)
        _tbox(slide, "🔑 Key Competitive Moves",
              _COL2_L + Inches(0.15), _CONTENT_TOP + Inches(0.15),
              _COL2_W - Inches(0.30), Inches(0.38),
              size=12, bold=True, color=_NAVY)

        _moves = [
            ("Celestica",  "SD6300 switch",              "Nov 2025"),
            ("Jabil",      "Cloud infra expansion",      "Q1 2026"),
            ("Sanmina",    "AI end-market push",         "ongoing"),
            ("Flex",       "Liquid cooling launch",      "FY2026"),
        ]
        for i, (company, desc, date) in enumerate(_moves):
            row_top = _CONTENT_TOP + Inches(0.65) + i * Inches(0.96)
            box = slide.shapes.add_textbox(
                _COL2_L + Inches(0.20), row_top,
                _COL2_W - Inches(0.35), Inches(0.86))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run(); r1.text = company + "  "
            r1.font.bold = True; r1.font.size = Pt(12)
            r1.font.color.rgb = _rgb(*_NAVY)
            r2 = p.add_run(); r2.text = desc
            r2.font.size = Pt(12); r2.font.color.rgb = _rgb(*_DARK_GRAY)
            p2 = tf.add_paragraph()
            r3 = p2.add_run(); r3.text = date
            r3.font.size = Pt(10); r3.font.italic = True
            r3.font.color.rgb = _rgb(*_MID_GRAY)

    elif num == 8:
        # Slide 9 — Risk Landscape: "Risk Priority" matrix
        panel_h = Inches(4.80)
        _rect(slide, _COL2_L, _CONTENT_TOP, _COL2_W, panel_h,
              fill=(254, 252, 252), line=_BLUE)
        _tbox(slide, "⚠️ Risk Priority",
              _COL2_L + Inches(0.15), _CONTENT_TOP + Inches(0.15),
              _COL2_W - Inches(0.30), Inches(0.38),
              size=12, bold=True, color=_NAVY)

        _risks = [
            (_RED,         "🔴 HIGH",  "Celestica/Jabil AI acceleration"),
            (_ORANGE,      "🟠 MED",   "Potential AI market oversupply"),
            (_GREEN_BADGE, "🟢 LOW",   "Benchmark competitive threat"),
        ]
        for i, (color, level, desc) in enumerate(_risks):
            row_top = _CONTENT_TOP + Inches(0.65) + i * Inches(1.30)
            # Colored left border strip
            _rect(slide,
                  _COL2_L + Inches(0.14), row_top + Inches(0.08),
                  Inches(0.06), Inches(0.95), fill=color)
            box = slide.shapes.add_textbox(
                _COL2_L + Inches(0.30), row_top + Inches(0.10),
                _COL2_W - Inches(0.45), Inches(0.95))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run(); r1.text = level + "  "
            r1.font.bold = True; r1.font.size = Pt(11)
            r1.font.color.rgb = _rgb(*color)
            p2 = tf.add_paragraph()
            r2 = p2.add_run(); r2.text = desc
            r2.font.size = Pt(11); r2.font.color.rgb = _rgb(*_DARK_GRAY)

    elif num == 10:
        # Slide 11 — Next Steps: "90-Day Execution Timeline"
        panel_h = Inches(4.80)
        _rect(slide, _COL2_L, _CONTENT_TOP, _COL2_W, panel_h,
              fill=_LIGHT_BLUE, line=_BLUE)
        _tbox(slide, "📅 Execution Timeline",
              _COL2_L + Inches(0.15), _CONTENT_TOP + Inches(0.15),
              _COL2_W - Inches(0.30), Inches(0.38),
              size=12, bold=True, color=_NAVY)

        _timeline = [
            (_BLUE,        "30d", "Assess AI market opportunities"),
            (_ORANGE,      "60d", "Strengthen hyperscaler partnerships"),
            (_GREEN_BADGE, "90d", "Increase CapEx in AI regions"),
        ]
        for i, (color, period, desc) in enumerate(_timeline):
            row_top = _CONTENT_TOP + Inches(0.72) + i * Inches(1.30)
            # Vertical connector line between circles (not after last)
            if i < len(_timeline) - 1:
                _rect(slide,
                      _COL2_L + Inches(0.32), row_top + Inches(0.46),
                      Inches(0.04), Inches(1.30), fill=(203, 213, 225))
            # Colored circle
            circ = slide.shapes.add_shape(
                9,  # oval
                _COL2_L + Inches(0.14), row_top + Inches(0.04),
                Inches(0.40), Inches(0.40))
            circ.fill.solid()
            circ.fill.fore_color.rgb = _rgb(*color)
            circ.line.fill.background()
            # Period label overlaid on circle
            _tbox(slide, period,
                  _COL2_L + Inches(0.14), row_top + Inches(0.07),
                  Inches(0.40), Inches(0.30),
                  size=9, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
            # Description text beside circle
            _tbox(slide, desc,
                  _COL2_L + Inches(0.65), row_top + Inches(0.08),
                  _COL2_W - Inches(0.80), Inches(0.36),
                  size=11, color=_DARK_GRAY)

    else:
        # Fallback: key-insight box — height sized to content
        insight_text = (insight or "")[:200]
        # ~44 chars fit per line at 12pt in the available width
        num_lines = max(1, (len(insight_text) + 43) // 44)
        box_h = Inches(min(0.48 + num_lines * 0.40 + 0.30,
                           1.2 if num_lines == 1 else 2.0))
        _rect(slide, _COL2_L, _CONTENT_TOP, _COL2_W, box_h,
              fill=_LIGHT_BLUE, line=_BLUE)
        _tbox(slide, "Key Insight",
              _COL2_L + Inches(0.15), _CONTENT_TOP + Inches(0.12),
              _COL2_W - Inches(0.3), Inches(0.34),
              size=11, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)
        _tbox(slide, insight_text,
              _COL2_L + Inches(0.15), _CONTENT_TOP + Inches(0.50),
              _COL2_W - Inches(0.3), box_h - Inches(0.60),
              size=11, italic=True, color=_DARK_GRAY)

    _footer_strip(slide, action)


# ── SLIDE ROUTER ─────────────────────────────────────────────────────────────

def _render_slide(prs, slide_data, exec_liner):
    """Route each slide number to its specialised renderer."""
    n = slide_data.get("slide_number", 0)
    if   n == 1: _add_exec_summary_slide(prs, slide_data, exec_liner)
    elif n == 3: _add_ranking_slide(prs, slide_data)
    elif n == 5: _add_financial_slide(prs, slide_data)
    elif n == 9: _add_flex_position_slide(prs, slide_data)
    else:        _add_default_slide(prs, slide_data)


# ═══════════════════════════════════════════════════════════════════════════
# PUBLIC ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════

def generate_powerpoint_report(company: Optional[str] = None) -> bytes:
    """
    Generate an AI-powered competitive intelligence PowerPoint report.

    Args:
        company: Company name (e.g. "Flex") or None / "all" for full comparison.

    Returns:
        PPTX file as bytes.

    Raises:
        ImportError:  python-pptx not installed.
        ValueError:   LLM returned unparseable JSON (message safe to show user).
        Exception:    Any other unexpected error during generation.
    """
    if not HAS_PPTX:
        raise ImportError(
            "python-pptx is required for PowerPoint export. "
            "Install with: pip install python-pptx"
        )

    # Normalise: treat "all" or "comparison" as no filter
    company_filter: Optional[str] = None
    if company and company.lower() not in ("all", "comparison"):
        company_filter = company

    # ── Step 1: RAG retrieval ─────────────────────────────────────────────
    logger.info("Building RAG context for PowerPoint report (company=%s)", company_filter or "all")
    rag_context = _build_rag_context(company_filter)

    # ── Step 2: LLM call ─────────────────────────────────────────────────
    logger.info("Calling LLM to generate slide content …")
    raw_response = _call_llm_for_report(rag_context, company_filter)

    # ── Step 3: Parse JSON ────────────────────────────────────────────────
    # ValueError is raised here with a user-safe message if parsing fails
    data = _parse_llm_response(raw_response)

    exec_liner = data.get("executive_summary_one_liner", "")
    slides     = data.get("slides", [])

    # ── Step 4: Render PPTX ───────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H

    # Cover slide
    _add_title_slide(prs, data)

    # 10 content slides
    for slide_spec in slides:
        _render_slide(prs, slide_spec, exec_liner)

    # Serialise
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    logger.info("PowerPoint report generated: %d slides", 1 + len(slides))
    return output.getvalue()
